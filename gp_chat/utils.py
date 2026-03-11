# utils.py:
import os
import sys
import yaml
import tempfile
import subprocess
import io
import hashlib
from importlib import resources
import streamlit as st
from . import config

# --- Library Imports ---

# Word (.docx) support
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PowerPoint (.pptx) text-only support (Linux compatible)
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

@st.cache_data
def load_prompts():
    """パッケージ内のprompts.yamlを読み込む"""
    try:
        with resources.open_text("gp_chat", "prompts.yaml") as f:
            yaml_data = yaml.safe_load(f)
            return yaml_data.get("prompts", {})
    except Exception as e:
        print(f"Warning: prompts.yaml load failed: {e}")
        return {}

def find_env_files(directory="env2"):
    """ディレクトリ内の.envファイルを検索"""
    if not os.path.isdir(directory):
        return []
    return [os.path.join(directory, f) for f in os.listdir(directory) if f.endswith(".env")]

def extract_text_from_docx(file_bytes):
    """docxファイルからテキストを抽出"""
    if not HAS_DOCX:
        return "[Error] python-docx library is not installed."

    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        return f"[Error parsing docx] {str(e)}"

def extract_text_from_pptx(file_bytes):
    """
    PPTXからテキストのみを抽出 (Linux/Cloud Run対応)
    """
    if not HAS_PPTX:
        return "[Error] python-pptx library is not installed."

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        full_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # シェイプ内のテキストを収集
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            if slide_text:
                full_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

        return "\n".join(full_text)
    except Exception as e:
        return f"[Error parsing pptx] {str(e)}"

def run_pylint_validation(canvas_code, canvas_index, prompts):
    """コードに対してpylintを実行し、分析プロンプトを生成"""
    if not canvas_code or canvas_code.strip() == "" or canvas_code.strip() == config.ACE_EDITOR_DEFAULT_CODE.strip():
        st.toast(config.UITexts.NO_CODE_TO_VALIDATE, icon="⚠️")
        return

    spinner_text = config.UITexts.VALIDATE_SPINNER_MULTI.format(i=canvas_index + 1) if st.session_state['multi_code_enabled'] else config.UITexts.VALIDATE_SPINNER_SINGLE
    with st.spinner(spinner_text):
        tmp_file_path = ""
        pylint_report = ""
        try:
            with tempfile.NamedTemporaryFile(mode='w+', suffix='.py', delete=False, encoding='utf-8') as tmp_file:
                tmp_file_path = tmp_file.name
                tmp_file.write(canvas_code.replace('\r\n', '\n'))
                tmp_file.flush()

            # pylint実行
            result = subprocess.run(
                [sys.executable, "-m", "pylint", tmp_file_path],
                capture_output=True, text=True, check=False
            )

            error_output = (result.stderr or "") + (result.stdout or "")
            if "syntax-error" in error_output.lower():
                st.toast(config.UITexts.PYLINT_SYNTAX_ERROR, icon="⚠️")
                return 

            issues = []
            if result.stdout:
                issues = [line for line in result.stdout.splitlines() if line.strip() and not line.startswith(('*', '-')) and 'Your code has been rated' not in line]

            if issues:
                cleaned_issues = [issue.replace(f'{tmp_file_path}:', 'Line ') for issue in issues]
                pylint_report = "\n".join(cleaned_issues)
        finally:
            if os.path.exists(tmp_file_path):
                os.remove(tmp_file_path)

    if not pylint_report.strip():
        st.sidebar.success(f"✅ Canvas-{canvas_index + 1}: pylint検証完了。問題なし。")
        return

    # Geminiへの分析依頼プロンプト
    validation_template = prompts.get("validation", {}).get("text", "以下はpylintのレポートです。解析してください:\n{pylint_report}\n\n対象コード:\n{code_for_prompt}")
    code_for_prompt = f"```python\n{canvas_code}\n```"
    validation_prompt = validation_template.format(code_for_prompt=code_for_prompt, pylint_report=pylint_report)

    system_message = st.session_state['messages'][0] if st.session_state['messages'] and st.session_state['messages'][0]["role"] == "system" else {"role": "system", "content": ""}
    st.session_state['special_generation_messages'] = [system_message, {"role": "user", "content": validation_prompt}]
    st.session_state['is_generating'] = True

def process_uploaded_files_for_gemini(uploaded_files):
    """StreamlitのアップロードファイルをGemini API用に変換"""
    from google.genai import types

    api_parts = []
    display_info = []

    for uploaded_file in uploaded_files:
        file_bytes = uploaded_file.getvalue()
        mime_type = uploaded_file.type
        filename = uploaded_file.name
        file_ext = os.path.splitext(filename)[1].lower()

        # 1. PDF & Images (Gemini Native Support)
        if mime_type == "application/pdf" or mime_type.startswith("image/"):
            api_parts.append(types.Part.from_bytes(data=file_bytes, mime_type=mime_type))
            display_info.append({"name": filename, "type": mime_type, "size": len(file_bytes)})

        # 2. Word (.docx) -> Text Extraction
        elif "wordprocessingml" in mime_type or filename.endswith(".docx"):
            text_content = extract_text_from_docx(file_bytes)
            prompt_text = f"\n\n[Attached Document: {filename}]\n{text_content}\n"
            api_parts.append(types.Part.from_text(text=prompt_text))
            display_info.append({"name": filename, "type": "docx (text)", "size": len(file_bytes)})

        # 3. PowerPoint (.pptx) -> Text Extraction (Linux Safe)
        elif file_ext in [".pptx", ".ppt"]:
            text_content = extract_text_from_pptx(file_bytes)
            prompt_text = f"\n\n[Attached Presentation: {filename}]\n(Note: Text content extracted. For visual analysis, please upload as PDF.)\n\n{text_content}\n"
            api_parts.append(types.Part.from_text(text=prompt_text))
            display_info.append({"name": filename, "type": "pptx (text)", "size": len(file_bytes)})

        # 4. Code & Text files
        elif mime_type.startswith("text/") or filename.endswith((".py", ".js", ".md", ".txt", ".json", ".csv", ".yaml", ".yml")):
            try:
                text_content = file_bytes.decode("utf-8")
                prompt_text = f"\n\n[Attached File: {filename}]\n```\n{text_content}\n```\n"
                api_parts.append(types.Part.from_text(text=prompt_text))
                display_info.append({"name": filename, "type": "text", "size": len(file_bytes)})
            except Exception:
                 st.warning(f"Could not decode text file: {filename}")
        else:
            st.warning(f"Unsupported file type: {filename}")

    return api_parts, display_info

def load_app_config():
    """パッケージ内のconfig.yamlを読み込む"""
    try:
        with resources.open_text("gp_chat", "config.yaml") as f:
            return yaml.safe_load(f)
    except Exception:
        return {}